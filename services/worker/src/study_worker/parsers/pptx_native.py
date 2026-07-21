"""Read-only native PPTX extraction using python-pptx plus OOXML metadata."""

from __future__ import annotations

import asyncio
import hashlib
from pathlib import Path
from typing import Any, cast

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide

from study_contracts import BlockType
from study_worker.parsers.normalize import (
    MetadataValue,
    RawArtifact,
    RawBlock,
    RawBoundingBox,
    RawDocument,
    RawPage,
)
from study_worker.parsers.ooxml import (
    OOXMLSecurityError,
    OOXMLSecurityPolicy,
    SlideFeatures,
    inspect_slide_features,
)
from study_worker.parsers.protocols import ParserCapability, ParseRequest
from study_worker.parsers.router import NativeParserError

PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_EMU_PER_PIXEL = 9_525


class PPTXNativeParser:
    def __init__(
        self,
        *,
        max_pages: int,
        max_pixels: int,
        ooxml_policy: OOXMLSecurityPolicy | None = None,
    ) -> None:
        if max_pages <= 0 or max_pixels <= 0:
            raise ValueError("PPTX parser limits must be positive")
        self._max_pages = max_pages
        self._max_pixels = max_pixels
        self._ooxml_policy = ooxml_policy or OOXMLSecurityPolicy()
        self._capability = ParserCapability(
            profile="native-v1",
            source_backend="pptx-native",
            source_version="1.0",
            media_types=frozenset({PPTX_MEDIA_TYPE}),
        )

    @property
    def capability(self) -> ParserCapability:
        return self._capability

    async def parse(self, request: ParseRequest) -> RawDocument:
        return await asyncio.to_thread(self.parse_sync, request)

    def parse_sync(self, request: ParseRequest) -> RawDocument:
        if request.media_type != PPTX_MEDIA_TYPE:
            raise NativeParserError("UNSUPPORTED_MEDIA_TYPE")
        _verify_input(request.input_path, request.document_sha256)
        try:
            features = inspect_slide_features(request.input_path, self._ooxml_policy)
        except OOXMLSecurityError as exc:
            raise NativeParserError(exc.code) from None
        try:
            presentation = Presentation(str(request.input_path))
        except Exception:
            raise NativeParserError("PPTX_PARSE_FAILED", retryable=True) from None
        total_page_count = len(presentation.slides)
        if total_page_count == 0:
            raise NativeParserError("PPTX_EMPTY")
        if total_page_count > self._max_pages:
            raise NativeParserError("PAGE_LIMIT_EXCEEDED")
        selected = _selected_ordinals(request.requested_pages, total_page_count)
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        if slide_width is None or slide_height is None:
            raise NativeParserError("PPTX_PAGE_SIZE_MISSING")
        width = max(1, round(int(slide_width) / _EMU_PER_PIXEL))
        height = max(1, round(int(slide_height) / _EMU_PER_PIXEL))
        if width * height > self._max_pixels:
            raise NativeParserError("PIXEL_LIMIT_EXCEEDED")
        request.output_dir.mkdir(mode=0o700, parents=True, exist_ok=True)
        pages = [
            self._parse_slide(
                presentation.slides[ordinal - 1],
                ordinal=ordinal,
                width=width,
                height=height,
                output_dir=request.output_dir,
                features=features.get(ordinal, _empty_features()),
            )
            for ordinal in selected
        ]
        return RawDocument(
            document_sha256=request.document_sha256,
            source_backend="pptx-native",
            source_version="1.0",
            total_page_count=total_page_count,
            pages=pages,
        )

    def _parse_slide(
        self,
        slide: Slide,
        *,
        ordinal: int,
        width: int,
        height: int,
        output_dir: Path,
        features: SlideFeatures,
    ) -> RawPage:
        drafts: list[
            tuple[
                float,
                float,
                BlockType,
                str,
                RawBoundingBox,
                dict[str, MetadataValue],
                RawArtifact | None,
            ]
        ] = []
        title_shape_id = getattr(slide.shapes.title, "shape_id", None)
        for shape in slide.shapes:
            dynamic_shape = cast(Any, shape)
            box = _shape_bbox(shape, width=width, height=height)
            top = box.top
            left = box.x0
            shape_name = str(getattr(shape, "name", "unnamed"))
            if bool(getattr(shape, "has_table", False)):
                table = dynamic_shape.table
                rows = [
                    [cell.text.replace("\r", " ").replace("\n", " ").strip() for cell in row.cells]
                    for row in table.rows
                ]
                drafts.append(
                    (
                        top,
                        left,
                        BlockType.TABLE,
                        "\n".join("\t".join(row) for row in rows).strip(),
                        box,
                        {
                            "shape_name": shape_name,
                            "row_count": len(rows),
                            "column_count": max((len(row) for row in rows), default=0),
                        },
                        None,
                    )
                )
                continue
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                image = dynamic_shape.image
                blob = bytes(image.blob)
                extension = str(image.ext).lower()
                relative_path = Path(
                    "assets", f"slide-{ordinal:06d}-picture-{shape.shape_id:06d}.{extension}"
                )
                artifact_path = output_dir / relative_path
                artifact_path.parent.mkdir(mode=0o700, parents=True, exist_ok=True)
                artifact_path.write_bytes(blob)
                pixel_width, pixel_height = image.size
                if int(pixel_width) * int(pixel_height) > self._max_pixels:
                    raise NativeParserError("PIXEL_LIMIT_EXCEEDED")
                artifact = RawArtifact(
                    relative_path=relative_path.as_posix(),
                    media_type=str(image.content_type),
                    sha256=hashlib.sha256(blob).hexdigest(),
                    size_bytes=len(blob),
                )
                drafts.append(
                    (
                        top,
                        left,
                        BlockType.IMAGE,
                        "",
                        box,
                        {
                            "shape_name": shape_name,
                            "pixel_width": int(pixel_width),
                            "pixel_height": int(pixel_height),
                        },
                        artifact,
                    )
                )
                continue
            if bool(getattr(shape, "has_text_frame", False)):
                text = str(getattr(shape, "text", "")).strip()
                if not text:
                    continue
                block_type = (
                    BlockType.TITLE
                    if getattr(shape, "shape_id", None) == title_shape_id
                    else BlockType.PARAGRAPH
                )
                drafts.append(
                    (
                        top,
                        left,
                        block_type,
                        text,
                        box,
                        {
                            "shape_name": shape_name,
                            "shape_type": int(getattr(shape, "shape_type", 0) or 0),
                        },
                        None,
                    )
                )

        for formula_index, formula_text in enumerate(features.omml_texts, start=1):
            drafts.append(
                (
                    float(height),
                    float(formula_index),
                    BlockType.FORMULA,
                    formula_text,
                    RawBoundingBox(x0=0, top=0, x1=1, bottom=1),
                    {"omml": True, "formula_index": formula_index},
                    None,
                )
            )
        if features.smartart_count:
            drafts.append(
                (
                    float(height),
                    float(width),
                    BlockType.IMAGE,
                    "",
                    RawBoundingBox(x0=0, top=0, x1=1, bottom=1),
                    {"smartart": True, "item_count": features.smartart_count},
                    None,
                )
            )
        for ole_index, ole in enumerate(features.ole_objects, start=1):
            drafts.append(
                (
                    float(height),
                    float(width + ole_index),
                    BlockType.IMAGE,
                    "",
                    RawBoundingBox(x0=0, top=0, x1=1, bottom=1),
                    {
                        "ole": True,
                        "ole_index": ole_index,
                        "name": ole.name,
                        "prog_id": ole.prog_id,
                    },
                    None,
                )
            )

        drafts.sort(key=lambda item: (item[0], item[1], item[2].value))
        blocks = [
            RawBlock(
                type=block_type,
                text=text,
                bbox=box,
                reading_order=reading_order,
                metadata=metadata,
                artifact=artifact,
            )
            for reading_order, (_, _, block_type, text, box, metadata, artifact) in enumerate(
                drafts
            )
        ]
        return RawPage(
            ordinal=ordinal,
            width=width,
            height=height,
            source_kind="slide",
            native_text_present=any(block.text for block in blocks),
            blocks=blocks,
            metadata={
                "external_relationship_count": features.external_relationship_count,
                "ole_count": len(features.ole_objects),
                "smartart_count": features.smartart_count,
                "omml_count": len(features.omml_texts),
            },
        )


def _verify_input(path: Path, expected_sha256: str) -> None:
    try:
        digest = hashlib.sha256(path.read_bytes()).hexdigest()
    except OSError:
        raise NativeParserError("PPTX_INPUT_UNREADABLE", retryable=True) from None
    if digest != expected_sha256:
        raise NativeParserError("INPUT_HASH_MISMATCH")


def _selected_ordinals(requested: tuple[int, ...], total: int) -> tuple[int, ...]:
    selected = requested or tuple(range(1, total + 1))
    if any(ordinal > total for ordinal in selected):
        raise NativeParserError("REQUESTED_PAGE_INVALID")
    return selected


def _shape_bbox(shape: object, *, width: int, height: int) -> RawBoundingBox:
    x0 = max(0.0, int(getattr(shape, "left", 0)) / _EMU_PER_PIXEL)
    top = max(0.0, int(getattr(shape, "top", 0)) / _EMU_PER_PIXEL)
    x1 = min(float(width), x0 + int(getattr(shape, "width", 0)) / _EMU_PER_PIXEL)
    bottom = min(float(height), top + int(getattr(shape, "height", 0)) / _EMU_PER_PIXEL)
    return RawBoundingBox(x0=x0, top=top, x1=max(x0, x1), bottom=max(top, bottom))


def _empty_features() -> SlideFeatures:
    return SlideFeatures(
        omml_texts=(),
        smartart_count=0,
        ole_objects=(),
        external_relationship_count=0,
    )
