"""Build deterministic PDF/PPTX fixtures without private course material."""

from __future__ import annotations

import base64
import hashlib
import zipfile
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

_FIXED_ZIP_TIME = (2026, 1, 1, 0, 0, 0)
_ONE_PIXEL_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
)


@dataclass(frozen=True, slots=True)
class DocumentFixtures:
    pdf: Path
    pptx: Path
    macro_pptx: Path


def build_documents(output_dir: Path) -> DocumentFixtures:
    output_dir.mkdir(parents=True, exist_ok=True)
    pdf = output_dir / "self-authored-native.pdf"
    pptx = output_dir / "self-authored-native.pptx"
    macro_pptx = output_dir / "self-authored-macro.pptx"
    pdf.write_bytes(_build_pdf_bytes())
    _build_pptx(pptx)
    _add_macro_part(pptx, macro_pptx)
    return DocumentFixtures(pdf=pdf, pptx=pptx, macro_pptx=macro_pptx)


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def rewrite_zip(
    source: Path,
    destination: Path,
    *,
    additions: dict[str, bytes] | None = None,
) -> None:
    entries: dict[str, bytes] = {}
    with zipfile.ZipFile(source) as archive:
        for info in archive.infolist():
            entries[info.filename] = archive.read(info)
    entries.update(additions or {})
    _write_canonical_zip(destination, entries)


def _build_pdf_bytes() -> bytes:
    page_one = b"\n".join(
        (
            b"BT /F1 28 Tf 72 730 Td (Operating Systems Review) Tj ET",
            b"BT /F1 24 Tf 72 700 Td (Processes use virtual memory and scheduling.) Tj ET",
            b"BT /F1 20 Tf 72 670 Td (Large body text remains paragraph text.) Tj ET",
            b"72 620 m 360 620 l 360 540 l 72 540 l h S",
            b"216 620 m 216 540 l S",
            b"72 580 m 360 580 l S",
            b"BT /F1 10 Tf 82 596 Td (Concept) Tj ET",
            b"BT /F1 10 Tf 226 596 Td (Definition) Tj ET",
            b"BT /F1 10 Tf 82 556 Td (Semaphore) Tj ET",
            b"BT /F1 10 Tf 226 556 Td (Synchronization primitive) Tj ET",
            b"q 96 0 0 48 72 430 cm /Im1 Do Q",
        )
    )
    page_two = b"q 240 0 0 160 180 300 cm /Im1 Do Q\n"
    image = b"\x20\x80\xe0"
    objects = (
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R 6 0 R] /Count 2 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 5 0 R >> /XObject << /Im1 7 0 R >> >> "
            b"/Contents 4 0 R >>"
        ),
        _pdf_stream(page_one),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /XObject << /Im1 7 0 R >> >> /Contents 8 0 R >>"
        ),
        (
            b"<< /Type /XObject /Subtype /Image /Width 1 /Height 1 "
            b"/ColorSpace /DeviceRGB /BitsPerComponent 8 /Length 3 >>\nstream\n"
            + image
            + b"\nendstream"
        ),
        _pdf_stream(page_two),
    )
    output = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for object_number, payload in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{object_number} 0 obj\n".encode())
        output.extend(payload)
        output.extend(b"\nendobj\n")
    xref_offset = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode())
    output.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode()
    )
    return bytes(output)


def _pdf_stream(payload: bytes) -> bytes:
    return f"<< /Length {len(payload)} >>\nstream\n".encode() + payload + b"endstream"


def _build_pptx(destination: Path) -> None:
    with TemporaryDirectory() as temporary_directory:
        temporary = Path(temporary_directory)
        image_path = temporary / "self-authored.png"
        image_path.write_bytes(_ONE_PIXEL_PNG)
        raw_path = temporary / "raw.pptx"

        presentation = Presentation()
        presentation.core_properties.author = "Self-authored fixture"
        presentation.core_properties.title = "Operating Systems Review"
        presentation.core_properties.created = datetime(2026, 1, 1)
        presentation.core_properties.modified = datetime(2026, 1, 1)
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Operating Systems Review"
        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5.6), Inches(0.8))
        text_box.text_frame.text = "Processes, memory, and synchronization"
        table = slide.shapes.add_table(
            2, 2, Inches(0.7), Inches(2.3), Inches(5.5), Inches(1.4)
        ).table
        table.cell(0, 0).text = "Concept"
        table.cell(0, 1).text = "Definition"
        table.cell(1, 0).text = "Semaphore"
        table.cell(1, 1).text = "Synchronization primitive"
        rectangle = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.2), Inches(2.2), Inches(0.8)
        )
        rectangle.text = "Scheduler"
        slide.shapes.add_picture(
            str(image_path), Inches(3.4), Inches(4.2), Inches(1.2), Inches(0.8)
        )
        presentation.save(str(raw_path))

        with zipfile.ZipFile(raw_path) as archive:
            entries = {info.filename: archive.read(info) for info in archive.infolist()}
        entries["ppt/slides/slide1.xml"] = _inject_slide_features(entries["ppt/slides/slide1.xml"])
        entries["ppt/slides/_rels/slide1.xml.rels"] = _inject_relationships(
            entries["ppt/slides/_rels/slide1.xml.rels"]
        )
        entries["[Content_Types].xml"] = _inject_content_types(entries["[Content_Types].xml"])
        entries["ppt/diagrams/data1.xml"] = (
            b'<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"/>'
        )
        entries["ppt/embeddings/oleObject1.bin"] = b"SELF_AUTHORED_OPAQUE_OLE_METADATA_ONLY"
        _write_canonical_zip(destination, entries)


def _inject_slide_features(payload: bytes) -> bytes:
    insertion = b"""
<p:graphicFrame>
  <p:nvGraphicFramePr>
    <p:cNvPr id="700" name="SelfAuthoredSmartArt"/>
    <p:cNvGraphicFramePr/><p:nvPr/>
  </p:nvGraphicFramePr>
  <p:xfrm><a:off x="0" y="0"/><a:ext cx="1" cy="1"/></p:xfrm>
  <a:graphic>
    <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram">
      <dgm:relIds
        xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram"
        r:dm="rIdSmartArt"/>
    </a:graphicData>
  </a:graphic>
</p:graphicFrame>
<p:graphicFrame>
  <p:nvGraphicFramePr>
    <p:cNvPr id="701" name="SelfAuthoredOLE"/>
    <p:cNvGraphicFramePr/><p:nvPr/>
  </p:nvGraphicFramePr>
  <p:xfrm><a:off x="0" y="0"/><a:ext cx="1" cy="1"/></p:xfrm>
  <a:graphic>
    <a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/ole">
      <p:oleObj name="MetadataOnly" progId="Package" r:id="rIdOle"/>
    </a:graphicData>
  </a:graphic>
</p:graphicFrame>
<m:oMath xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math"><m:r><m:t>x+y</m:t></m:r></m:oMath>
"""
    return payload.replace(b"</p:spTree>", insertion + b"</p:spTree>")


def _inject_relationships(payload: bytes) -> bytes:
    insertion = b"""
<Relationship Id="rIdSmartArt"
  Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData"
  Target="../diagrams/data1.xml"/>
<Relationship Id="rIdOle"
  Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject"
  Target="../embeddings/oleObject1.bin"/>
<Relationship Id="rIdExternal"
  Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"
  Target="https://example.invalid/self-authored" TargetMode="External"/>
"""
    return payload.replace(b"</Relationships>", insertion + b"</Relationships>")


def _inject_content_types(payload: bytes) -> bytes:
    insertion = b"""
<Default Extension="bin" ContentType="application/vnd.openxmlformats-officedocument.oleObject"/>
<Override PartName="/ppt/diagrams/data1.xml"
  ContentType="application/vnd.openxmlformats-officedocument.drawingml.diagramData+xml"/>
"""
    return payload.replace(b"</Types>", insertion + b"</Types>")


def _add_macro_part(source: Path, destination: Path) -> None:
    rewrite_zip(
        source,
        destination,
        additions={"ppt/vbaProject.bin": b"SELF_AUTHORED_MACRO_SENTINEL_DO_NOT_EXECUTE"},
    )


def _write_canonical_zip(destination: Path, entries: dict[str, bytes]) -> None:
    with zipfile.ZipFile(destination, "w") as archive:
        for name in sorted(entries):
            info = zipfile.ZipInfo(name, date_time=_FIXED_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            archive.writestr(
                info,
                entries[name],
                compress_type=zipfile.ZIP_DEFLATED,
                compresslevel=9,
            )
